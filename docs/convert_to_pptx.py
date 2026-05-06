#!/usr/bin/env python3
"""
Convert KDM Consortium HTML slides to PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 bold=False, color=None, align=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to the slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = RgbColor(*color)
    
    return txBox

def add_bullet_text(slide, left, top, width, height, bullets, font_size=14):
    """Add bullet points to the slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"✓ {bullet}"
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = 'Calibri'
        p.font.color.rgb = RgbColor(255, 255, 255)
    
    return txBox

def create_ppt_from_html():
    """Create PowerPoint from HTML slides"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio (1280x720 equivalent)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    DARK_BLUE = RgbColor(13, 31, 51)      # #0d1f33
    MEDIUM_BLUE = RgbColor(30, 58, 95)    # #1e3a5f
    GOLD = RgbColor(245, 158, 11)         # #f59e0b
    WHITE = RgbColor(255, 255, 255)
    LIGHT_GRAY = RgbColor(148, 163, 184)  # #94a3b8
    GREEN = RgbColor(34, 197, 94)       # #22c55e
    
    # Slide 1: Welcome
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Set background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MEDIUM_BLUE
    
    # Add gradient overlay effect (simulated with shapes)
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.fill.fore_color.brightness = 0.2
    shape.line.fill.background()
    
    # Logo
    add_text_box(slide1, 0.5, 0.3, 4, 0.5, "KDM CONSORTIUM", 
                font_size=24, bold=True, color=(245, 158, 11))
    
    # Pill badge
    pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   Inches(0.5), Inches(0.9), Inches(2.5), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RgbColor(245, 158, 11)
    pill.fill.fore_color.brightness = 0.8
    pill.line.fill.background()
    
    add_text_box(slide1, 0.6, 0.95, 2.3, 0.3, "SELECTIVE NETWORK",
                font_size=14, color=(245, 158, 11))
    
    # Main title
    add_text_box(slide1, 0.5, 1.5, 9, 1.5,
                "Join an Exclusive Network of Government Contracting Experts",
                font_size=44, bold=True, color=(255, 255, 255))
    
    # Subtitle
    add_text_box(slide1, 0.5, 3.0, 8, 0.5,
                "12-50 Expert Companies. One Mission: Winning Together.",
                font_size=24, color=(148, 163, 184))
    
    # Three cards
    card_data = [
        {
            "title": "🎯 Selective Model",
            "bullets": ["12-50 curated members only", "Hand-picked by capability fit", "High-touch, not mass market"]
        },
        {
            "title": "🏆 Five Pillars",
            "bullets": ["U.S. Manufacturing", "Critical Minerals", "Defense Contracting", "Access to Capital", "Opportunity Zones"]
        },
        {
            "title": "💰 Membership",
            "bullets": ["$1,250/month standard", "$650/month promo", "Valid through May 2026"]
        }
    ]
    
    for i, card in enumerate(card_data):
        x_pos = 0.5 + i * 4.0
        
        # Card background
        card_shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(x_pos), Inches(4.0), Inches(3.5), Inches(2.5))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RgbColor(255, 255, 255)
        card_shape.fill.fore_color.brightness = -0.9  # 10% opacity effect
        card_shape.line.color.rgb = RgbColor(255, 255, 255)
        card_shape.line.color.brightness = -0.85
        
        # Card title
        add_text_box(slide1, x_pos + 0.2, 4.2, 3.1, 0.4, card["title"],
                    font_size=20, bold=True, color=(245, 158, 11))
        
        # Card bullets
        add_bullet_text(slide1, x_pos + 0.2, 4.7, 3.1, 1.8, card["bullets"], font_size=14)
    
    # Slide number
    add_text_box(slide1, 11.5, 6.8, 1.5, 0.3, "1 / 6",
                font_size=16, color=(255, 255, 255), align=PP_ALIGN.RIGHT)
    
    # Save the presentation
    output_path = r"C:\Users\Buyer\Documents\CascadeProjects\KDMAssociates\svp-platform\docs\KDM-Consortium-Slides.pptx"
    prs.save(output_path)
    print(f"PowerPoint created: {output_path}")
    
    return output_path

if __name__ == "__main__":
    try:
        output_file = create_ppt_from_html()
        print(f"\n✅ Successfully created: {output_file}")
        print("\nNote: This creates Slide 1 only. For all 6 slides, the script can be extended.")
    except Exception as e:
        print(f"\n❌ Error: {e}")
        print("\nMake sure python-pptx is installed:")
        print("  pip install python-pptx beautifulsoup4")
